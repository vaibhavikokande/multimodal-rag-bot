"""
Document Processing Engine
Handles PDF, DOCX, PPTX, TXT, CSV, XLSX with OCR support
"""
import os
import io
import time
import asyncio
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from loguru import logger

from app.core.config import settings


class DocumentProcessor:
    """Main document processing orchestrator"""

    SUPPORTED_TYPES = {
        "pdf": "_process_pdf",
        "docx": "_process_docx",
        "doc": "_process_docx",
        "pptx": "_process_pptx",
        "ppt": "_process_pptx",
        "txt": "_process_txt",
        "csv": "_process_csv",
        "xlsx": "_process_xlsx",
        "xls": "_process_xlsx",
        "png": "_process_image",
        "jpg": "_process_image",
        "jpeg": "_process_image",
        "gif": "_process_image",
        "bmp": "_process_image",
        "tiff": "_process_image",
        "webp": "_process_image",
        "mp4": "_process_video",
        "avi": "_process_video",
        "mov": "_process_video",
        "mkv": "_process_video",
        "mp3": "_process_audio",
        "wav": "_process_audio",
        "m4a": "_process_audio",
        "ogg": "_process_audio",
    }

    async def process(self, file_path: str, file_type: str) -> Dict:
        """Process a document and return extracted content"""
        start_time = time.time()
        file_type = file_type.lower().strip(".")

        if file_type not in self.SUPPORTED_TYPES:
            raise ValueError(f"Unsupported file type: {file_type}")

        method_name = self.SUPPORTED_TYPES[file_type]
        method = getattr(self, method_name)

        try:
            result = await method(file_path)
            result["processing_time"] = time.time() - start_time
            result["file_type"] = file_type
            return result
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            raise

    async def _process_pdf(self, file_path: str) -> Dict:
        """Extract text from PDF with OCR fallback"""
        import pdfplumber
        from pdf2image import convert_from_path

        text_content = []
        tables = []
        metadata = {}
        page_count = 0
        has_tables = False

        try:
            with pdfplumber.open(file_path) as pdf:
                metadata = pdf.metadata or {}
                page_count = len(pdf.pages)

                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text() or ""

                    # OCR fallback for scanned PDFs
                    if len(page_text.strip()) < 50:
                        page_text = await self._ocr_page(file_path, page_num - 1)

                    if page_text:
                        text_content.append(f"[Page {page_num}]\n{page_text}")

                    # Extract tables
                    page_tables = page.extract_tables()
                    if page_tables:
                        has_tables = True
                        for table in page_tables:
                            if table:
                                table_text = self._table_to_text(table)
                                tables.append({
                                    "page": page_num,
                                    "content": table_text
                                })

            return {
                "text": "\n\n".join(text_content),
                "tables": tables,
                "metadata": metadata,
                "page_count": page_count,
                "has_tables": has_tables,
                "has_images": True,
            }
        except Exception as e:
            logger.warning(f"pdfplumber failed, trying OCR: {e}")
            return await self._full_ocr_pdf(file_path)

    async def _full_ocr_pdf(self, file_path: str) -> Dict:
        """Full OCR for scanned PDFs"""
        from pdf2image import convert_from_path
        import pytesseract

        images = convert_from_path(file_path, dpi=300)
        text_pages = []

        for i, image in enumerate(images, 1):
            text = pytesseract.image_to_string(image, lang="eng")
            if text.strip():
                text_pages.append(f"[Page {i}]\n{text}")

        return {
            "text": "\n\n".join(text_pages),
            "tables": [],
            "metadata": {},
            "page_count": len(images),
            "has_tables": False,
            "has_images": True,
        }

    async def _ocr_page(self, file_path: str, page_index: int) -> str:
        """OCR a single PDF page"""
        try:
            from pdf2image import convert_from_path
            import pytesseract
            images = convert_from_path(
                file_path, dpi=300,
                first_page=page_index + 1,
                last_page=page_index + 1
            )
            if images:
                return pytesseract.image_to_string(images[0])
        except Exception as e:
            logger.warning(f"OCR failed for page {page_index}: {e}")
        return ""

    async def _process_docx(self, file_path: str) -> Dict:
        """Extract text from DOCX"""
        from docx import Document as DocxDocument
        from docx.oxml.ns import qn

        doc = DocxDocument(file_path)
        paragraphs = []
        tables = []

        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name if para.style else "Normal"
                if "Heading" in style:
                    paragraphs.append(f"\n## {para.text}\n")
                else:
                    paragraphs.append(para.text)

        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                table_text = "\n".join(rows)
                tables.append({"content": table_text})
                paragraphs.append(f"\n[TABLE]\n{table_text}\n")

        metadata = {
            "author": doc.core_properties.author,
            "title": doc.core_properties.title,
            "created": str(doc.core_properties.created),
        }

        return {
            "text": "\n".join(paragraphs),
            "tables": tables,
            "metadata": metadata,
            "page_count": None,
            "has_tables": len(tables) > 0,
            "has_images": False,
        }

    async def _process_pptx(self, file_path: str) -> Dict:
        """Extract text from PPTX"""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation(file_path)
        slides_text = []
        has_images = False

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = [f"[Slide {slide_num}]"]

            if slide.name:
                slide_content.append(f"Title: {slide.name}")

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_content.append(text)

                if shape.shape_type == 13:  # Picture
                    has_images = True
                    slide_content.append("[IMAGE: Visual content on this slide]")

            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_content.append(f"Notes: {notes_text}")

            slides_text.append("\n".join(slide_content))

        metadata = {
            "slide_count": len(prs.slides),
            "author": prs.core_properties.author,
            "title": prs.core_properties.title,
        }

        return {
            "text": "\n\n".join(slides_text),
            "tables": [],
            "metadata": metadata,
            "page_count": len(prs.slides),
            "has_tables": False,
            "has_images": has_images,
        }

    async def _process_txt(self, file_path: str) -> Dict:
        """Process plain text files"""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return {
            "text": text,
            "tables": [],
            "metadata": {},
            "page_count": None,
            "has_tables": False,
            "has_images": False,
        }

    async def _process_csv(self, file_path: str) -> Dict:
        """Process CSV files"""
        import pandas as pd
        df = pd.read_csv(file_path)
        text_parts = []

        # Schema description
        text_parts.append(f"CSV File with {len(df)} rows and {len(df.columns)} columns.")
        text_parts.append(f"Columns: {', '.join(df.columns.tolist())}")

        # Sample data
        text_parts.append("\nFirst 5 rows:")
        text_parts.append(df.head(5).to_string())

        # Statistical summary
        text_parts.append("\nStatistical Summary:")
        text_parts.append(df.describe().to_string())

        # Full data as text
        text_parts.append("\nFull Data:")
        text_parts.append(df.to_string())

        return {
            "text": "\n".join(text_parts),
            "tables": [{"content": df.to_string()}],
            "metadata": {"rows": len(df), "columns": len(df.columns)},
            "page_count": None,
            "has_tables": True,
            "has_images": False,
        }

    async def _process_xlsx(self, file_path: str) -> Dict:
        """Process Excel files"""
        import pandas as pd
        xl = pd.ExcelFile(file_path)
        text_parts = []
        tables = []

        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            text_parts.append(f"\n[Sheet: {sheet_name}]")
            text_parts.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
            text_parts.append(df.to_string())
            tables.append({"sheet": sheet_name, "content": df.to_string()})

        return {
            "text": "\n".join(text_parts),
            "tables": tables,
            "metadata": {"sheets": xl.sheet_names},
            "page_count": len(xl.sheet_names),
            "has_tables": True,
            "has_images": False,
        }

    async def _process_image(self, file_path: str) -> Dict:
        """Process images with OCR and AI captioning"""
        from PIL import Image
        import pytesseract

        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)

        metadata = {
            "width": image.width,
            "height": image.height,
            "mode": image.mode,
            "format": image.format,
        }

        return {
            "text": text or "[Image content - no extractable text]",
            "tables": [],
            "metadata": metadata,
            "page_count": 1,
            "has_tables": False,
            "has_images": True,
        }

    async def _process_video(self, file_path: str) -> Dict:
        """Process video files - extract audio and transcribe"""
        from app.services.video_processor.processor import VideoProcessor
        vp = VideoProcessor()
        return await vp.process(file_path)

    async def _process_audio(self, file_path: str) -> Dict:
        """Process audio files with Whisper"""
        from app.services.video_processor.processor import VideoProcessor
        vp = VideoProcessor()
        transcript = await vp.transcribe_audio(file_path)
        return {
            "text": transcript,
            "tables": [],
            "metadata": {},
            "page_count": None,
            "has_tables": False,
            "has_images": False,
            "transcript": transcript,
        }

    def _table_to_text(self, table: List) -> str:
        """Convert table array to readable text"""
        rows = []
        for row in table:
            if row:
                cleaned = [str(cell).strip() if cell else "" for cell in row]
                rows.append(" | ".join(cleaned))
        return "\n".join(rows)
